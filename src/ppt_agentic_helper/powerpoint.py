from __future__ import annotations

from datetime import date, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .models import GanttPlan, GanttStyle, GanttTask, ManagedGanttSlide, TaskVisualStyle
from .planning import metadata_text_to_slide, slide_to_metadata_text

METADATA_SHAPE_NAME = "PPT_AGENTIC_HELPER_METADATA"


def create_or_update_presentation(
    *,
    plan: GanttPlan,
    presentation_path: str,
    slide_title: str,
    style: GanttStyle | None = None,
    include_metadata: bool = True,
    output_path: str | None = None,
) -> str:
    path = Path(presentation_path)
    if path.exists():
        presentation = Presentation(str(path))
    else:
        presentation = Presentation()

    slide = find_slide_by_title(presentation, slide_title)
    if slide is None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    rebuild_slide(
        slide,
        plan=plan,
        slide_title=slide_title,
        style=style or GanttStyle(),
        include_metadata=include_metadata,
    )

    final_path = Path(output_path) if output_path else path
    if not final_path.suffix:
        final_path = final_path.with_suffix(".pptx")
    final_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(final_path))
    return str(final_path.resolve())


def load_plan_from_presentation(*, presentation_path: str, slide_title: str) -> GanttPlan:
    return load_managed_slide(presentation_path=presentation_path, slide_title=slide_title).plan


def load_style_from_presentation(*, presentation_path: str, slide_title: str) -> GanttStyle:
    return load_managed_slide(presentation_path=presentation_path, slide_title=slide_title).style


def load_managed_slide(*, presentation_path: str, slide_title: str) -> ManagedGanttSlide:
    path = Path(presentation_path)
    if not path.exists():
        raise FileNotFoundError(f"Presentation not found: {presentation_path}")

    presentation = Presentation(str(path))
    slide = find_slide_by_title(presentation, slide_title)
    if slide is None:
        raise RuntimeError(f'Could not find a slide titled "{slide_title}"')

    for shape in slide.shapes:
        if shape.name == METADATA_SHAPE_NAME and hasattr(shape, "text"):
            metadata = metadata_text_to_slide(shape.text)
            if metadata.slide_title != slide_title:
                metadata.slide_title = slide_title
            return metadata

    raise RuntimeError(
        f'Slide "{slide_title}" was found, but it does not contain helper metadata to support updates'
    )


def bootstrap_style_from_presentation(
    *,
    presentation_path: str,
    slide_title: str | None = None,
    slide_index: int | None = None,
    theme_name: str = "bootstrapped",
) -> GanttStyle:
    path = Path(presentation_path)
    if not path.exists():
        raise FileNotFoundError(f"Presentation not found: {presentation_path}")

    if slide_title is None and slide_index is None:
        raise ValueError("slide_title or slide_index is required")

    presentation = Presentation(str(path))
    slide = get_slide(presentation, slide_title=slide_title, slide_index=slide_index)
    return bootstrap_style_from_slide(slide, theme_name=theme_name)


def bootstrap_style_from_slide(slide: Slide, *, theme_name: str = "bootstrapped") -> GanttStyle:
    style = GanttStyle(theme_name=theme_name)
    slide_width = slide.part.package.presentation_part.presentation.slide_width
    slide_height = slide.part.package.presentation_part.presentation.slide_height

    background_shape = detect_background_shape(slide, slide_width=slide_width, slide_height=slide_height)
    if background_shape is not None:
        style.slide_background_color = shape_fill_hex(background_shape, style.slide_background_color)

    title_shape = detect_title_shape(slide)
    if title_shape is not None:
        style.title_color = shape_text_hex(title_shape, style.title_color)

    summary_shape = detect_summary_shape(slide, title_shape=title_shape)
    if summary_shape is not None:
        style.summary_color = shape_text_hex(summary_shape, style.summary_color)

    header_shapes = detect_header_shapes(slide, title_shape=title_shape, summary_shape=summary_shape)
    if header_shapes:
        reference = header_shapes[0]
        style.header_fill_color = shape_fill_hex(reference, style.header_fill_color)
        style.header_border_color = shape_line_hex(reference, style.header_border_color)
        style.header_text_color = shape_text_hex(reference, style.header_text_color)

    row_shapes = detect_row_background_shapes(slide, header_shapes=header_shapes)
    if row_shapes:
        style.row_fill_color = shape_fill_hex(row_shapes[0], style.row_fill_color)
        style.grid_border_color = shape_line_hex(row_shapes[0], style.grid_border_color)
        if len(row_shapes) > 1:
            style.alternate_row_fill_color = shape_fill_hex(row_shapes[1], style.alternate_row_fill_color)

    task_rows = detect_task_row_textboxes(slide, header_shapes=header_shapes)
    if task_rows:
        style.task_name_color = first_paragraph_text_hex(task_rows[0], style.task_name_color)
        style.owner_color = nth_paragraph_text_hex(task_rows[0], 1, style.owner_color)

    bar_shapes = detect_task_bar_shapes(slide, header_shapes=header_shapes)
    matched_rows = match_rows_to_bars(task_rows, bar_shapes)
    status_visuals: dict[str, list[TaskVisualStyle]] = {}
    task_styles: dict[str, TaskVisualStyle] = {}
    for row_shape, bar_shape in matched_rows:
        task_name = extract_task_name(row_shape)
        if not task_name:
            continue

        visual = TaskVisualStyle(
            fill_color=shape_fill_hex(bar_shape, "#9CA3AF"),
            border_color=shape_line_hex(bar_shape, "#6B7280"),
            text_color=shape_text_hex(bar_shape, "#FFFFFF"),
        )
        task_styles[task_name] = visual

        status = infer_status_from_text(row_shape.text)
        if status:
            status_visuals.setdefault(status, []).append(visual)

    style.task_styles = task_styles
    status_styles: dict[str, TaskVisualStyle] = {}
    for status_name, visuals in status_visuals.items():
        if len(visuals) < 2:
            continue
        keys = {task_visual_style_key(visual) for visual in visuals}
        if len(keys) == 1:
            status_styles[status_name] = visuals[0]
    for status_name, default_visual in GanttStyle().status_styles.items():
        style.status_styles[status_name] = status_styles.get(status_name, default_visual)

    return style


def get_slide(
    presentation: Presentation,
    *,
    slide_title: str | None = None,
    slide_index: int | None = None,
) -> Slide:
    if slide_title is not None:
        slide = find_slide_by_title(presentation, slide_title)
        if slide is None:
            raise RuntimeError(f'Could not find a slide titled "{slide_title}"')
        return slide

    if slide_index is None:
        raise ValueError("slide_index is required when slide_title is not provided")

    zero_based_index = slide_index - 1
    if zero_based_index < 0 or zero_based_index >= len(presentation.slides):
        raise IndexError(f"slide_index {slide_index} is out of range")
    return presentation.slides[zero_based_index]


def find_slide_by_title(presentation: Presentation, slide_title: str) -> Slide | None:
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape.text.strip() == slide_title:
                return slide
    return None


def rebuild_slide(
    slide: Slide,
    *,
    plan: GanttPlan,
    slide_title: str,
    style: GanttStyle,
    include_metadata: bool,
) -> None:
    clear_slide(slide)
    add_background(slide, style)
    add_title(slide, slide_title, style)
    add_summary(slide, plan, style)
    draw_gantt_grid(slide, plan, style)
    if include_metadata:
        add_metadata(slide, ManagedGanttSlide(slide_title=slide_title, plan=plan, style=style))


def clear_slide(slide: Slide) -> None:
    shape_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        shape_tree.remove(shape._element)


def add_background(slide: Slide, style: GanttStyle) -> None:
    slide_width = slide.part.package.presentation_part.presentation.slide_width
    slide_height = slide.part.package.presentation_part.presentation.slide_height
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = color_from_hex(style.slide_background_color)
    background.line.fill.background()


def add_title(slide: Slide, slide_title: str, style: GanttStyle) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8.5), Inches(0.5))
    title_box.name = "PPT_AGENTIC_HELPER_TITLE"
    text_frame = title_box.text_frame
    text_frame.text = slide_title
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True
    paragraph.font.color.rgb = color_from_hex(style.title_color)


def add_summary(slide: Slide, plan: GanttPlan, style: GanttStyle) -> None:
    summary_text = plan.summary or "Generated from notes."
    if plan.assumptions:
        assumptions = " | ".join(plan.assumptions[:3])
        summary_text = f"{summary_text} Assumptions: {assumptions}"

    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.0), Inches(0.55))
    summary_box.name = "PPT_AGENTIC_HELPER_SUMMARY"
    text_frame = summary_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = summary_text
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(10)
    paragraph.font.color.rgb = color_from_hex(style.summary_color)


def draw_gantt_grid(slide: Slide, plan: GanttPlan, style: GanttStyle) -> None:
    top = Inches(1.45)
    left = Inches(0.5)
    table_width = Inches(3.55)
    chart_left = left + table_width
    total_width = Inches(12.3)
    chart_width = total_width - table_width
    bottom_margin = Inches(0.45)
    row_header_height = Inches(0.45)
    slide_height = slide.part.package.presentation_part.presentation.slide_height
    usable_height = slide_height - top - bottom_margin
    row_height = max(int((usable_height - row_header_height) / max(len(plan.tasks), 1)), int(Inches(0.32)))

    weeks = timeline_columns(plan)
    if not weeks:
        weeks = [week_start(plan.start_date)]
    week_width = int(chart_width / len(weeks))

    draw_headers(slide, left, top, table_width, chart_left, row_header_height, weeks, week_width, style)

    task_top = top + row_header_height
    for row_index, task in enumerate(plan.tasks):
        row_y = task_top + row_index * row_height
        add_row_background(slide, left, row_y, total_width, row_height, row_index, style)
        add_task_cells(slide, task, left, row_y, table_width, row_height, style)
        add_task_bar(slide, task, chart_left, row_y, row_height, weeks, week_width, style)


def draw_headers(
    slide: Slide,
    left: int,
    top: int,
    table_width: int,
    chart_left: int,
    row_header_height: int,
    weeks: list[date],
    week_width: int,
    style: GanttStyle,
) -> None:
    header_fill = color_from_hex(style.header_fill_color)
    border = color_from_hex(style.header_border_color)

    task_header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, table_width, row_header_height
    )
    task_header.fill.solid()
    task_header.fill.fore_color.rgb = header_fill
    task_header.line.color.rgb = border
    task_header.text = "Task / Owner / Status"
    format_cell_text(task_header, 11, bold=True, color=color_from_hex(style.header_text_color))

    for index, week in enumerate(weeks):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            chart_left + index * week_width,
            top,
            week_width,
            row_header_height,
        )
        box.fill.solid()
        box.fill.fore_color.rgb = header_fill
        box.line.color.rgb = border
        box.text = week.strftime("%b %d")
        format_cell_text(
            box,
            9,
            align=PP_ALIGN.CENTER,
            bold=True,
            color=color_from_hex(style.header_text_color),
        )


def add_row_background(
    slide: Slide, left: int, top: int, width: int, height: int, row_index: int, style: GanttStyle
) -> None:
    fill = (
        color_from_hex(style.row_fill_color)
        if row_index % 2 == 0
        else color_from_hex(style.alternate_row_fill_color)
    )
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = fill
    background.line.color.rgb = color_from_hex(style.grid_border_color)


def add_task_cells(
    slide: Slide, task: GanttTask, left: int, top: int, width: int, height: int, style: GanttStyle
) -> None:
    box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.03), width - Inches(0.16), height)
    text_frame = box.text_frame
    text_frame.word_wrap = True

    p1 = text_frame.paragraphs[0]
    p1.text = task.name
    p1.font.bold = True
    p1.font.size = Pt(11)
    p1.font.color.rgb = color_from_hex(style.task_name_color)

    owner = task.owner or "Unassigned"
    p2 = text_frame.add_paragraph()
    p2.text = owner
    p2.font.size = Pt(9)
    p2.font.color.rgb = color_from_hex(style.owner_color)

    p3 = text_frame.add_paragraph()
    p3.text = task.status.replace("_", " ").title()
    p3.font.size = Pt(9)
    p3.font.color.rgb = color_from_hex(resolve_task_style(style, task).border_color or "#6B7280")


def add_task_bar(
    slide: Slide,
    task: GanttTask,
    chart_left: int,
    top: int,
    row_height: int,
    weeks: list[date],
    week_width: int,
    style: GanttStyle,
) -> None:
    start_index = max(0, (task.start_date - weeks[0]).days // 7)
    end_index = max(start_index, (task.end_date - weeks[0]).days // 7)
    bar_left = chart_left + start_index * week_width + Inches(0.05)
    bar_top = top + Inches(0.09)
    bar_width = max((end_index - start_index + 1) * week_width - Inches(0.1), Inches(0.18))
    bar_height = max(row_height - Inches(0.18), Inches(0.14))

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, bar_left, bar_top, bar_width, bar_height)
    bar.fill.solid()
    task_style = resolve_task_style(style, task)
    bar.fill.fore_color.rgb = color_from_hex(task_style.fill_color or "#9CA3AF")
    bar.line.color.rgb = color_from_hex(task_style.border_color or "#6B7280")

    label = task.owner or task.name
    bar.text = label[:24]
    format_cell_text(
        bar,
        9,
        align=PP_ALIGN.CENTER,
        color=color_from_hex(task_style.text_color or "#FFFFFF"),
        bold=True,
    )


def add_metadata(slide: Slide, managed_slide: ManagedGanttSlide) -> None:
    metadata = slide.shapes.add_textbox(Inches(0.05), Inches(7.0), Inches(0.2), Inches(0.2))
    metadata.name = METADATA_SHAPE_NAME
    metadata.text = slide_to_metadata_text(managed_slide)

    text_frame = metadata.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(1)
    paragraph.font.color.rgb = RGBColor(255, 255, 255)


def format_cell_text(
    shape,
    font_size: int,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    color: RGBColor = RGBColor(31, 41, 55),
    bold: bool = False,
) -> None:
    text_frame = shape.text_frame
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold


def week_start(value: date) -> date:
    return value - timedelta(days=value.weekday())


def timeline_columns(plan: GanttPlan) -> list[date]:
    start = week_start(plan.start_date)
    end = week_start(plan.end_date)
    total_weeks = ((end - start).days // 7) + 1
    return [start + timedelta(weeks=offset) for offset in range(total_weeks)]


def resolve_task_style(style: GanttStyle, task: GanttTask) -> TaskVisualStyle:
    if task.name in style.task_styles:
        task_style = style.task_styles[task.name]
        status_style = style.status_styles.get(task.status, TaskVisualStyle())
        return TaskVisualStyle(
            fill_color=task_style.fill_color or status_style.fill_color,
            border_color=task_style.border_color or status_style.border_color,
            text_color=task_style.text_color or status_style.text_color,
        )
    return style.status_styles.get(task.status, TaskVisualStyle())


def color_from_hex(value: str) -> RGBColor:
    hex_value = value.strip().removeprefix("#")
    if len(hex_value) != 6:
        raise ValueError(f"Expected 6-digit hex color, got: {value}")
    return RGBColor.from_string(hex_value.upper())


def detect_background_shape(slide: Slide, *, slide_width: int, slide_height: int):
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "fill"):
            continue
        area = getattr(shape, "width", 0) * getattr(shape, "height", 0)
        if area < slide_width * slide_height * 0.7:
            continue
        candidates.append(shape)
    if not candidates:
        return None
    return max(candidates, key=lambda shape: shape.width * shape.height)


def detect_title_shape(slide: Slide):
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
    if not text_shapes:
        return None
    return max(text_shapes, key=lambda shape: (largest_font_size(shape), -shape.top))


def detect_summary_shape(slide: Slide, *, title_shape):
    if title_shape is None:
        return None
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape is not title_shape
        and shape.text.strip()
        and shape.top > title_shape.top
        and shape.top < Inches(1.6)
    ]
    if not candidates:
        return None
    return min(candidates, key=lambda shape: shape.top)


def detect_header_shapes(slide: Slide, *, title_shape, summary_shape) -> list:
    min_top = None
    shapes = []
    cutoff = Inches(1.25)
    if summary_shape is not None:
        cutoff = max(cutoff, summary_shape.top + summary_shape.height - Inches(0.1))
    elif title_shape is not None:
        cutoff = max(cutoff, title_shape.top + title_shape.height)

    for shape in slide.shapes:
        if getattr(shape, "name", "") == METADATA_SHAPE_NAME:
            continue
        if not hasattr(shape, "fill") or getattr(shape, "top", 0) < cutoff:
            continue
        if getattr(shape, "height", 0) < Inches(0.2):
            continue
        if getattr(shape, "top", 0) > Inches(2.2):
            continue
        min_top = shape.top if min_top is None else min(min_top, shape.top)
        shapes.append(shape)

    if min_top is None:
        return []

    return sorted([shape for shape in shapes if abs(shape.top - min_top) <= Inches(0.15)], key=lambda shape: shape.left)


def detect_row_background_shapes(slide: Slide, *, header_shapes: list) -> list:
    header_bottom = max((shape.top + shape.height for shape in header_shapes), default=Inches(1.8))
    rows = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "fill")
        and getattr(shape, "top", 0) >= header_bottom
        and getattr(shape, "width", 0) >= Inches(7.0)
        and getattr(shape, "height", 0) >= Inches(0.2)
        and not getattr(shape, "has_text_frame", False)
    ]
    rows = sorted(rows, key=lambda shape: shape.top)
    distinct = []
    seen_tops = []
    for shape in rows:
        if any(abs(shape.top - top) <= Inches(0.05) for top in seen_tops):
            continue
        distinct.append(shape)
        seen_tops.append(shape.top)
    return distinct[:2]


def detect_task_row_textboxes(slide: Slide, *, header_shapes: list) -> list:
    header_bottom = max((shape.top + shape.height for shape in header_shapes), default=Inches(1.8))
    rows = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip()
        and getattr(shape, "left", 0) < Inches(4.3)
        and getattr(shape, "top", 0) >= header_bottom
        and getattr(shape, "height", 0) >= Inches(0.2)
    ]
    return sorted(rows, key=lambda shape: shape.top)


def detect_task_bar_shapes(slide: Slide, *, header_shapes: list) -> list:
    header_bottom = max((shape.top + shape.height for shape in header_shapes), default=Inches(1.8))
    bars = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "fill")
        and getattr(shape, "left", 0) > Inches(3.8)
        and getattr(shape, "top", 0) >= header_bottom
        and getattr(shape, "width", 0) >= Inches(0.2)
        and getattr(shape, "height", 0) >= Inches(0.12)
    ]
    return sorted(bars, key=lambda shape: (shape.top, shape.left))


def match_rows_to_bars(task_rows: list, bar_shapes: list) -> list[tuple[object, object]]:
    matches: list[tuple[object, object]] = []
    remaining_bars = list(bar_shapes)
    for row_shape in task_rows:
        row_center = row_shape.top + row_shape.height / 2
        compatible = [
            bar for bar in remaining_bars if abs((bar.top + bar.height / 2) - row_center) <= max(row_shape.height, Inches(0.3))
        ]
        if not compatible:
            continue
        best_bar = min(compatible, key=lambda bar: abs((bar.top + bar.height / 2) - row_center))
        matches.append((row_shape, best_bar))
        remaining_bars.remove(best_bar)
    return matches


def largest_font_size(shape) -> float:
    sizes = []
    if not getattr(shape, "has_text_frame", False):
        return 0
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size:
            sizes.append(float(paragraph.font.size))
        for run in paragraph.runs:
            if run.font.size:
                sizes.append(float(run.font.size))
    return max(sizes, default=0)


def extract_task_name(shape) -> str | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            return text
    return None


def infer_status_from_text(text: str) -> str | None:
    normalized = text.lower().replace("-", " ").replace("_", " ")
    if "not started" in normalized:
        return "not_started"
    if "in progress" in normalized:
        return "in_progress"
    if "blocked" in normalized:
        return "blocked"
    if "complete" in normalized or "completed" in normalized:
        return "complete"
    return None


def shape_fill_hex(shape, default: str) -> str:
    try:
        rgb = shape.fill.fore_color.rgb
        if rgb is None:
            return default
        return rgb_to_hex(rgb)
    except Exception:
        return default


def shape_line_hex(shape, default: str) -> str:
    try:
        rgb = shape.line.color.rgb
        if rgb is None:
            return default
        return rgb_to_hex(rgb)
    except Exception:
        return default


def shape_text_hex(shape, default: str) -> str:
    return first_paragraph_text_hex(shape, default)


def first_paragraph_text_hex(shape, default: str) -> str:
    return nth_paragraph_text_hex(shape, 0, default)


def nth_paragraph_text_hex(shape, index: int, default: str) -> str:
    try:
        paragraph = shape.text_frame.paragraphs[index]
    except Exception:
        return default

    try:
        if paragraph.font.color.rgb is not None:
            return rgb_to_hex(paragraph.font.color.rgb)
    except Exception:
        pass

    for run in paragraph.runs:
        try:
            if run.font.color.rgb is not None:
                return rgb_to_hex(run.font.color.rgb)
        except Exception:
            continue
    return default


def rgb_to_hex(rgb) -> str:
    return f"#{str(rgb).upper()}"


def task_visual_style_key(style: TaskVisualStyle) -> tuple[str | None, str | None, str | None]:
    return (style.fill_color, style.border_color, style.text_color)


def status_color(status: str) -> RGBColor:
    return {
        "not_started": RGBColor(107, 114, 128),
        "in_progress": RGBColor(37, 99, 235),
        "blocked": RGBColor(185, 28, 28),
        "complete": RGBColor(5, 150, 105),
    }.get(status, RGBColor(107, 114, 128))


def status_fill(status: str) -> RGBColor:
    return {
        "not_started": RGBColor(156, 163, 175),
        "in_progress": RGBColor(59, 130, 246),
        "blocked": RGBColor(239, 68, 68),
        "complete": RGBColor(16, 185, 129),
    }.get(status, RGBColor(156, 163, 175))
