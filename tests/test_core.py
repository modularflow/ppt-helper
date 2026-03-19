from datetime import date
from pathlib import Path

from pptx import Presentation

from ppt_agentic_helper.demo import run_capability_demo
from ppt_agentic_helper.models import GanttPlan, GanttTask
from ppt_agentic_helper.openai_chat import extract_json_object
from ppt_agentic_helper.planning import metadata_text_to_slide, plan_to_metadata_text
from ppt_agentic_helper.powerpoint import (
    bootstrap_style_from_presentation,
    create_or_update_presentation,
    load_managed_slide,
    load_plan_from_presentation,
    timeline_columns,
)
from ppt_agentic_helper.sample_data import sample_style


def sample_plan() -> GanttPlan:
    return GanttPlan(
        title="Migration Plan",
        summary="Coordinate discovery, build, and rollout.",
        assumptions=["Dates are estimated from meeting notes."],
        tasks=[
            GanttTask(
                name="Discovery",
                start_date=date(2026, 3, 2),
                end_date=date(2026, 3, 13),
                owner="PMO",
                status="complete",
            ),
            GanttTask(
                name="Build",
                start_date=date(2026, 3, 16),
                end_date=date(2026, 4, 3),
                owner="Engineering",
                status="in_progress",
                dependencies=["Discovery"],
            ),
        ],
    )


def test_extract_json_object_from_fenced_block() -> None:
    raw = '```json\n{"title":"Demo","tasks":[{"name":"Task","start_date":"2026-03-01","end_date":"2026-03-02"}],"assumptions":[]}\n```'
    assert extract_json_object(raw).startswith("{")


def test_metadata_round_trip() -> None:
    plan = sample_plan()
    restored = metadata_text_to_slide(plan_to_metadata_text(plan))
    assert restored.plan == plan


def test_timeline_columns_are_weekly() -> None:
    plan = sample_plan()
    weeks = timeline_columns(plan)
    assert weeks[0] == date(2026, 3, 2)
    assert weeks[-1] == date(2026, 3, 30)


def test_create_and_reload_presentation(tmp_path) -> None:
    plan = sample_plan()
    ppt_path = tmp_path / "timeline.pptx"

    create_or_update_presentation(
        plan=plan,
        presentation_path=str(ppt_path),
        slide_title="Migration Gantt",
    )

    restored = load_plan_from_presentation(
        presentation_path=str(ppt_path),
        slide_title="Migration Gantt",
    )

    assert restored == plan


def test_update_reuses_same_slide(tmp_path) -> None:
    original = sample_plan()
    revised = GanttPlan(
        title="Migration Plan",
        summary="Updated summary.",
        assumptions=["Adjusted after risk review."],
        tasks=[
            GanttTask(
                name="Discovery",
                start_date=date(2026, 3, 2),
                end_date=date(2026, 3, 13),
                owner="PMO",
                status="complete",
            ),
            GanttTask(
                name="Build",
                start_date=date(2026, 3, 16),
                end_date=date(2026, 4, 10),
                owner="Engineering",
                status="blocked",
                notes="Waiting on environment access.",
                dependencies=["Discovery"],
            ),
        ],
    )
    ppt_path = tmp_path / "timeline.pptx"

    create_or_update_presentation(
        plan=original,
        presentation_path=str(ppt_path),
        slide_title="Migration Gantt",
    )
    create_or_update_presentation(
        plan=revised,
        presentation_path=str(ppt_path),
        slide_title="Migration Gantt",
    )

    presentation = Presentation(str(ppt_path))
    assert len(presentation.slides) == 1

    restored = load_plan_from_presentation(
        presentation_path=str(ppt_path),
        slide_title="Migration Gantt",
    )
    assert restored == revised


def test_style_is_persisted_in_metadata(tmp_path) -> None:
    ppt_path = tmp_path / "styled.pptx"
    style = sample_style()

    create_or_update_presentation(
        plan=sample_plan(),
        presentation_path=str(ppt_path),
        slide_title="Styled Gantt",
        style=style,
    )

    managed_slide = load_managed_slide(
        presentation_path=str(ppt_path),
        slide_title="Styled Gantt",
    )
    assert managed_slide.style == style


def test_bootstrap_style_from_existing_slide(tmp_path) -> None:
    ppt_path = tmp_path / "reference.pptx"
    source_style = sample_style()

    create_or_update_presentation(
        plan=sample_plan(),
        presentation_path=str(ppt_path),
        slide_title="Reference Gantt",
        style=source_style,
        include_metadata=False,
    )

    inferred = bootstrap_style_from_presentation(
        presentation_path=str(ppt_path),
        slide_title="Reference Gantt",
        theme_name="inferred",
    )

    assert inferred.theme_name == "inferred"
    assert inferred.slide_background_color == source_style.slide_background_color
    assert inferred.header_fill_color == source_style.header_fill_color
    assert inferred.task_styles["Build"].fill_color == source_style.task_styles["Build"].fill_color


def test_capability_demo_outputs_all_artifacts(tmp_path) -> None:
    result = run_capability_demo(str(tmp_path / "demo"))

    expected_keys = {
        "output_dir",
        "source_presentation",
        "created_presentation",
        "updated_presentation",
        "source_style_json",
        "bootstrapped_style_json",
        "created_metadata_json",
        "updated_metadata_json",
        "creation_notes",
        "update_notes",
    }
    assert expected_keys.issubset(result)
    assert Path(result["created_presentation"]).exists()
    assert Path(result["updated_presentation"]).exists()
    Presentation(result["created_presentation"])
    Presentation(result["updated_presentation"])
    for key in expected_keys - {"output_dir"}:
        assert Path(result[key]).exists()
